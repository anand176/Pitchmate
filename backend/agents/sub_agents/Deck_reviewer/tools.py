"""
Deck Creator tools — build a pitch deck .pptx using python-pptx.
"""

import json
import os
import re
from datetime import datetime
from core.config import config

_ARTIFACTS_DIR = getattr(config, "artifacts_root_dir", "./artifacts")


def _sanitize_filename(name: str) -> str:
    """Return a safe filename stem (no path, no extension)."""
    stem = re.sub(r"[^\w\s-]", "", (name or "pitch_deck").strip())[:80]
    return stem or "pitch_deck"


def create_pitch_deck_pptx(content_json: str) -> str:
    """
    Create a PowerPoint pitch deck (.pptx) from structured JSON content and save it to the artifacts directory.

    Call this with a JSON string containing company_name, tagline, and a list of slides.
    Each slide has "title" and "bullets" (list of strings). The first slide is always the title slide;
    subsequent slides are title + bullet content.

    Args:
        content_json: JSON string with structure:
            {
              "company_name": "Startup Name",
              "tagline": "One-line tagline",
              "slides": [
                {"title": "The Problem", "bullets": ["Bullet 1", "Bullet 2"]},
                {"title": "Our Solution", "bullets": ["..."]},
                {"title": "Market", "bullets": ["..."]},
                ...
              ]
            }
            Standard pitch flow: Problem, Solution, Market (TAM/SAM/SOM), Product, Business Model,
            Traction, Competition, Team, Financials, The Ask, Thank You / Contact.

    Returns:
        Message with the file path to the created .pptx, or an error message.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return (
            "Deck creation failed: python-pptx is not installed. "
            "Install with: pip install python-pptx"
        )

    try:
        data = json.loads(content_json)
    except json.JSONDecodeError as e:
        return f"Invalid JSON for deck content: {e}. Provide a valid JSON with company_name, tagline, and slides (list of {{title, bullets}})."

    company_name = (data.get("company_name") or "Pitch Deck").strip()
    tagline = (data.get("tagline") or "").strip()
    slides_spec = data.get("slides")
    if not isinstance(slides_spec, list):
        return "JSON must include 'slides': a list of objects with 'title' and 'bullets' (list of strings)."

    os.makedirs(_ARTIFACTS_DIR, exist_ok=True)
    safe_name = _sanitize_filename(company_name)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"pitch_deck_{safe_name}_{timestamp}.pptx"
    filepath = os.path.join(_ARTIFACTS_DIR, filename)

    try:
        prs = Presentation()
        # Title slide (layout index 0: title + subtitle)
        title_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_layout)
        title_slide.shapes.title.text = company_name
        if tagline and len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = tagline

        # Content slides (layout index 1: title and body)
        content_layout = prs.slide_layouts[1]
        for spec in slides_spec:
            if not isinstance(spec, dict):
                continue
            title = (spec.get("title") or "Slide").strip()
            bullets = spec.get("bullets")
            if not isinstance(bullets, list):
                bullets = [str(spec.get("body", ""))] if spec.get("body") else []
            bullet_strs = [str(b).strip() for b in bullets if b]

            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = title
            if bullet_strs:
                body = slide.placeholders[1]
                tf = body.text_frame
                if hasattr(tf, "clear"):
                    tf.clear()
                for i, line in enumerate(bullet_strs):
                    if i == 0:
                        p = tf.paragraphs[0]
                        p.text = line
                        p.font.size = Pt(14)
                    else:
                        p = tf.add_paragraph()
                        p.text = line
                        p.level = 0
                        p.font.size = Pt(14)

        prs.save(filepath)
        abs_path = os.path.abspath(filepath)
        return f"Pitch deck created: {abs_path}"
    except Exception as e:
        return f"Deck creation failed: {str(e)}"
